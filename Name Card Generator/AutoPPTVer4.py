import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
import pandas as pd
import os
import sys


# 取得當前腳本所在的目錄
#script_dir = os.path.dirname(os.path.abspath(__file__))
script_dir = os.path.dirname(sys.argv[0])

# 指定PowerPoint文件的路徑
pptx_file_path = 'CBN_NAMECARD.pptx'
#resource_path = os.path.join(script_dir, 'AutoPPT')

# 構建完整的相對路徑
pptx_file_path = os.path.join(script_dir, pptx_file_path)
#pptx_file_path = os.path.join(resource_path, pptx_file_path)


''' Name Card PPT 資料生成'''
def modify_powerpoint(pptx_file_path, *replacement_sets):
    
    # 讀取PowerPoint文件
    presentation = Presentation(pptx_file_path)

    # 遍歷所有幻燈片
    for slide in presentation.slides:
        # 遍歷每個形狀
        for shape in slide.shapes:
            # 如果形狀是文本框，且包含目標文本，則進行替換
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for replacement_set in replacement_sets:
                            for target_text, replacement_text in replacement_set.items():
                                if target_text in run.text:
                                    run.text = run.text.replace(target_text, replacement_text)

    # 保存修改後的PowerPoint文件
    modified_pptx_file_path = pptx_file_path.replace('.pptx', '_modified.pptx')
    presentation.save(modified_pptx_file_path)

    print(f'Modified PowerPoint saved to: {modified_pptx_file_path}')

def close_window():
    global Name,JobTitle,TelN,MobileN,Mail,NameC,JobTitleC

    Name = name.get()
    JobTitle = jobtitle.get()
    TelN = telN.get()
    MobileN = mobileN.get()
    Mail = mail.get()
    NameC = nameC.get()
    JobTitleC = jobtitleC.get()

    # 定義多個替換字典
    replacement_set_1 = {'Name': Name}
    replacement_set_2 = {'JobTitle': JobTitle}
    replacement_set_3 = {'名字': NameC}
    replacement_set_4 = {'職稱': JobTitleC}
    replacement_set_5 = {'1234': TelN}
    replacement_set_6 = {'行動電話': MobileN}
    replacement_set_7 = {'MAIL_IN': Mail}
    
    # 呼叫函數進行多次修改
    modify_powerpoint(pptx_file_path, replacement_set_1, replacement_set_2, replacement_set_3,
                  replacement_set_4, replacement_set_5, replacement_set_6, replacement_set_7, 
                  )
    
    root.destroy()

# 使用者輸入的東西都算是字串
root = tk.Tk()
root.title('請輸入相關資訊')
root.geometry('500x500')

Name_Label = tk.Label(root, text="請輸入英文名 \n Ex:Mark Yu")
Name_Label.pack()
name = tk.Entry(root)
name.pack()

NameC_Label = tk.Label(root, text="請輸入中文名")
NameC_Label.pack()
nameC = tk.Entry(root)
nameC.pack()

JobTitle_Label = tk.Label(root, text="請輸入英文職稱")
JobTitle_Label.pack()
jobtitle = tk.Entry(root)
jobtitle.pack()

JobTitleC_Label = tk.Label(root, text="請輸入中文職稱")
JobTitleC_Label.pack()
jobtitleC = tk.Entry(root)
jobtitleC.pack()

#TelN_Label = tk.Label(root, text="請輸入公司電話(含分機) \n Ex:+886 3 560 0066 #2916 :")
TelN_Label = tk.Label(root, text="請輸入公司分機 \n Ex:2916 :")
TelN_Label.pack()
telN = tk.Entry(root)
telN.pack()

MobileN_Label = tk.Label(root, text="請輸入手機 \n Ex:+886 919 353 737 :")
MobileN_Label.pack()
mobileN = tk.Entry(root)
mobileN.pack()

Mail_Label = tk.Label(root, text="請輸入信箱 \n Ex:mark_yu@compalbn.com :")
Mail_Label.pack()
mail = tk.Entry(root)
mail.pack()

Enter_Button = tk.Button(root, text="開始生成名片", command=close_window)
Enter_Button.pack()

result_label = tk.Label(root, text="")
result_label.pack()

root.mainloop()


''' PPT轉換成JPG '''
import win32com.client
import os
import sys

path = os.path.abspath(os.path.dirname(sys.argv[0]))

# 獲取默認目錄下所有文件
filenames = os.listdir(path)

def ppt2png(ppt, filename):
    # 開啟PPT程序
    ppt_app = win32com.client.Dispatch('PowerPoint.Application')
    # 在後台開起ppt文件
    ppt = ppt_app.Presentations.Open(ppt)
    # 將開啟的PPT另存為圖片。17為ppt轉jpg,18為轉png，32為ppt轉pdf。
    ppt.SaveAs(filename, 17)
    print("已轉成JPG檔!!!")
    # 退出PPT程序
    ppt_app.Quit()

#循环讀取文件名稱
for ppt in filenames:
    #判断文件名稱，篩選出副檔名為ppt或pptx的文件
    if ppt.endswith('ppt') or ppt.endswith('pptx'):
        #將獲取的文件名字進行分割，提取出文件名，比如1024.ppt，按照下面的規則base=1024，ext=ppt
        base, ext = ppt.split('.') 
        #轉換成圖片後存放路徑
        filename = path + '/' + base + '.jpg'
        #識別出ppt文件後，拼接成ppt文件存放路徑
        ppt = path + '/' + ppt
        #執行ppt2png方法，轉換ppt
        ppt2png(ppt, filename)
    else:
        continue

print('開始生成QR code')

''' QR code 生成'''
# 載入需要的套件
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.webdriver.common.action_chains import ActionChains
from pynput.keyboard import Key, Controller
import pyautogui
import time
import os
import wget

# 開啟瀏覽器視窗(Chrome)
# 方法一：執行前需開啟chromedriver.exe且與執行檔在同一個工作目錄
driver = webdriver.Chrome()

#script_dir = os.path.dirname(os.path.abspath(__file__))
script_dir = os.path.dirname(sys.argv[0])

driver.get("https://me-qr.com/qr-code-generator/image?bannerid=5803082635&utm_source=google&utm_medium=cpc&utm_campaign=12760339093&utm_adgroupid=150040768077&utm_content=661296049630&audience=150040768077&keyword=%E7%85%A7%E7%89%87%20%E8%BD%89%20qr%20code%20%E7%B7%9A%E4%B8%8A&utm_target=&device=c&gclid=Cj0KCQiAwP6sBhDAARIsAPfK_wa3lH7oOnGAw6joGVpZI2bYOsAPQ21tOuC0GPpA_vUrUVJ0sigF8RIaAhXxEALw_wcB")
time.sleep(3)

driver.maximize_window()
# driver.fullscreen_window( )
time.sleep(5)

driver.find_element(By.CLASS_NAME,'filepond--wrapper').click()
file_input1 = driver.find_element(By.NAME, 'gallery')
time.sleep(10)

file_path1 = script_dir + "\CBN_NAMECARD_modified\投影片1.jpg"
file_input1.send_keys(file_path1)
time.sleep(10)

ele = driver.find_element(By.XPATH,'//*[@id="root"]/div[2]/div/section/div/div[1]')
ele.location_once_scrolled_into_view
time.sleep(2)

driver.find_element(By.XPATH,'//*[@id="secondStepClick"]').click()
time.sleep(5)

driver.find_element(By.XPATH, '//*[@id="downloadQrCode"]').click()
time.sleep(5)
driver.close()
driver.quit()
#time.sleep(2)

print("請自行去個人電腦的下載資料夾取得生成之QR code")
print("程式執行完畢，可以關閉程式")

'''
print('轉換第二張投影片!!')


# 執行第二張投影片轉QR code

driver = webdriver.Chrome()
driver.get("https://me-qr.com/qr-code-generator/image?bannerid=5803082635&utm_source=google&utm_medium=cpc&utm_campaign=12760339093&utm_adgroupid=150040768077&utm_content=661296049630&audience=150040768077&keyword=%E7%85%A7%E7%89%87%20%E8%BD%89%20qr%20code%20%E7%B7%9A%E4%B8%8A&utm_target=&device=c&gclid=Cj0KCQiAwP6sBhDAARIsAPfK_wa3lH7oOnGAw6joGVpZI2bYOsAPQ21tOuC0GPpA_vUrUVJ0sigF8RIaAhXxEALw_wcB")
time.sleep(3)

driver.maximize_window()
#driver.fullscreen_window( )
time.sleep(5)

driver.find_element(By.CLASS_NAME,'filepond--wrapper').click()
file_input2 = driver.find_element(By.NAME, 'gallery')
time.sleep(10)

file_path2 = script_dir + "\CBN_NAMECARD_modified\投影片2.jpg"
file_input2.send_keys(file_path2)
time.sleep(10)

ele = driver.find_element(By.XPATH,'//*[@id="root"]/div[2]/div/section/div/div[1]')
ele.location_once_scrolled_into_view
time.sleep(2)

driver.find_element(By.XPATH,'//*[@id="secondStepClick"]').click()
time.sleep(5)

driver.find_element(By.XPATH, '//*[@id="downloadQrCode"]').click()
time.sleep(5)
driver.close()
driver.quit()

print("請自行去個人電腦的下載資料夾取得生成之QR code")
print("程式執行完畢，可以關閉程式")
'''


'''
rootend = tk.Tk()
rootend.title('程式結束提醒')
rootend.geometry('300x300')

End_Label = tk.Label(rootend, text="請自行去個人電腦的下載資料夾取得生成之QR code")
End_Label.pack()

def end_close_window():
    rootend.destroy()

End_Button = tk.Button(rootend, text="結束程式", command=end_close_window)
End_Button.pack()

rootend.mainloop()
'''